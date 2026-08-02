import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

# PowerPointオブジェクトの作成
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ブランクスライドレイアウト
blank_slide_layout = prs.slide_layouts[6]

# カラーパレット定義
COLOR_BG_DARK = RGBColor(15, 23, 42)       # Slate 900 (Dark Navy)
COLOR_BG_LIGHT = RGBColor(248, 250, 252)   # Slate 50 (Light Grayish Blue)
COLOR_CARD_BG = RGBColor(255, 255, 255)    # White
COLOR_CARD_BORDER = RGBColor(226, 232, 240)# Slate 200

COLOR_PRIMARY = RGBColor(37, 99, 235)      # Blue 600
COLOR_PRIMARY_DARK = RGBColor(30, 58, 138) # Blue 900
COLOR_ACCENT_BLUE = RGBColor(14, 165, 233) # Sky 500
COLOR_ACCENT_AMBER = RGBColor(245, 158, 11)# Amber 500
COLOR_ACCENT_EMERALD = RGBColor(16, 185, 129)# Emerald 500

COLOR_TEXT_MAIN = RGBColor(30, 41, 59)     # Slate 800
COLOR_TEXT_MUTED = RGBColor(100, 116, 139) # Slate 500
COLOR_TEXT_WHITE = RGBColor(255, 255, 255) # White

FONT_NAME = "Meiryo"

def set_slide_bg(slide, color):
    """スライド背景色を設定"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header(slide, title_text, category_text="YK特化WMS"):
    """標準ヘッダーを追加"""
    # ヘッダー背景バー
    header_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
    header_box.fill.solid()
    header_box.fill.fore_color.rgb = COLOR_PRIMARY_DARK
    header_box.line.color.rgb = COLOR_PRIMARY_DARK

    # カテゴリタグ
    cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.12), Inches(10), Inches(0.3))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.name = FONT_NAME
    p_cat.font.size = Pt(11)
    p_cat.font.bold = True
    p_cat.font.color.rgb = COLOR_ACCENT_BLUE

    # タイトル
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.6))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = FONT_NAME
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_TEXT_WHITE

def add_card(slide, left, top, width, height, bg_color=COLOR_CARD_BG, border_color=COLOR_CARD_BORDER):
    """角丸カード背景を追加"""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(1)
    else:
        card.line.fill.background()
    return card

def add_image_with_border(slide, img_path, left, top, max_w, max_h):
    """アスペクト比を維持して画像枠線付きで挿入"""
    if not os.path.exists(img_path):
        # 画像が存在しない場合のフォールバック枠
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, max_w, max_h)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(241, 245, 249)
        box.line.color.rgb = RGBColor(203, 213, 225)
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Image Not Found: {os.path.basename(img_path)}"
        p.alignment = PP_ALIGN.CENTER
        p.font.name = FONT_NAME
        p.font.color.rgb = COLOR_TEXT_MUTED
        return

    im = Image.open(img_path)
    im_w, im_h = im.size
    aspect = im_w / im_h

    # 表示サイズの計算
    target_w = max_w
    target_h = max_w / aspect

    if target_h > max_h:
        target_h = max_h
        target_w = max_h * aspect

    # センタリングオフセット
    offset_x = (max_w - target_w) / 2
    offset_y = (max_h - target_h) / 2

    final_left = left + offset_x
    final_top = top + offset_y

    # 背後のシャドウ影枠
    shadow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, final_left + Inches(0.05), final_top + Inches(0.05), target_w, target_h)
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = RGBColor(203, 213, 225)
    shadow.line.fill.background()

    # 画像本体
    slide.shapes.add_picture(img_path, final_left, final_top, width=target_w, height=target_h)

    # 枠線
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, final_left, final_top, target_w, target_h)
    border.fill.background()
    border.line.color.rgb = RGBColor(148, 163, 184)
    border.line.width = Pt(1)


# ==========================================
# Slide 1: 表紙
# ==========================================
slide1 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(slide1, COLOR_BG_DARK)

# 装飾背景アクセント
bg_acc = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
bg_acc.fill.solid()
bg_acc.fill.fore_color.rgb = COLOR_PRIMARY
bg_acc.line.fill.background()

# メインタイトルテキスト
tb_t1 = slide1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11), Inches(1.2))
tf1 = tb_t1.text_frame
tf1.word_wrap = True
p1 = tf1.paragraphs[0]
p1.text = "YK特化WMS"
p1.font.name = FONT_NAME
p1.font.size = Pt(44)
p1.font.bold = True
p1.font.color.rgb = COLOR_TEXT_WHITE

# サブタイトル
tb_sub = slide1.shapes.add_textbox(Inches(1.2), Inches(3.0), Inches(11), Inches(1.0))
tf_sub = tb_sub.text_frame
tf_sub.word_wrap = True
p_sub = tf_sub.paragraphs[0]
p_sub.text = "次世代型 倉庫管理システム（WMS） システム概要 ＆ 機能詳細プレゼンテーション"
p_sub.font.name = FONT_NAME
p_sub.font.size = Pt(22)
p_sub.font.bold = True
p_sub.font.color.rgb = COLOR_ACCENT_BLUE

# 特徴バッジ
badges = ["バラ箱（残数）管理", "自動倉庫選定＆自動引当", "高度運賃最適化エンジン", "ワンストップ配送管理"]
for i, badge in enumerate(badges):
    bx = Inches(1.2 + i * 2.8)
    by = Inches(4.3)
    card_b = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, bx, by, Inches(2.6), Inches(0.6))
    card_b.fill.solid()
    card_b.fill.fore_color.rgb = RGBColor(30, 41, 59)
    card_b.line.color.rgb = COLOR_PRIMARY
    card_b.line.width = Pt(1.5)
    tf_b = card_b.text_frame
    p_b = tf_b.paragraphs[0]
    p_b.text = f"✔ {badge}"
    p_b.alignment = PP_ALIGN.CENTER
    p_b.font.name = FONT_NAME
    p_b.font.size = Pt(11)
    p_b.font.bold = True
    p_b.font.color.rgb = COLOR_TEXT_WHITE

# フッター
tb_ft = slide1.shapes.add_textbox(Inches(1.2), Inches(6.2), Inches(10), Inches(0.5))
p_ft = tb_ft.text_frame.paragraphs[0]
p_ft.text = "システム開発部 | 概要仕様および機能カタログ"
p_ft.font.name = FONT_NAME
p_ft.font.size = Pt(12)
p_ft.font.color.rgb = COLOR_TEXT_MUTED


# ==========================================
# Slide 2: WMSとしての強みと4大コア特徴
# ==========================================
slide2 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(slide2, COLOR_BG_LIGHT)
add_header(slide2, "WMSとしての強みと4大コア特徴", "SYSTEM CONCEPT & FEATURES")

features = [
    {
        "num": "01",
        "title": "バラ箱（残数）＆ FIFO管理",
        "color": COLOR_PRIMARY,
        "desc": "未開封ケースだけでなく、開封後のバラ箱の残存部数をリアルタイム追跡。FIFO（先入先出）× バラ箱優先引当により、現場の端数在庫を無駄なく回転させます。"
    },
    {
        "num": "02",
        "title": "最安倉庫の自動選定・出荷引当",
        "color": COLOR_ACCENT_BLUE,
        "desc": "出荷指示取り込み時に、各倉庫の在庫状況と配送運賃を瞬時に解析。最安で出荷可能な倉庫を自動割り振ることで物流コストを自動最適化します。"
    },
    {
        "num": "03",
        "title": "高度な運賃比較・シミュレーション",
        "color": COLOR_ACCENT_AMBER,
        "desc": "複数運送会社の運賃表・距離・重量・エリア・倉庫割増率・個別運賃を網羅。画面上でキャリア別の最安料金を可視化し比較・調整できます。"
    },
    {
        "num": "04",
        "title": "ワンストップ現場オペレーション",
        "color": COLOR_ACCENT_EMERALD,
        "desc": "Excel指示読込から倉庫確定、運賃シミュレーション、ヤマトB2/佐川e飛伝などの送り状CSV出力・ラベル発行まで一気通貫で完結します。"
    }
]

for i, feat in enumerate(features):
    col = i % 2
    row = i // 2
    left = Inches(0.8 + col * 5.9)
    top = Inches(1.5 + row * 2.7)
    
    add_card(slide2, left, top, Inches(5.6), Inches(2.4))
    
    # ナンバー丸アイコン
    circle = slide2.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.3), top + Inches(0.3), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = feat["color"]
    circle.line.fill.background()
    p_c = circle.text_frame.paragraphs[0]
    p_c.text = feat["num"]
    p_c.alignment = PP_ALIGN.CENTER
    p_c.font.name = FONT_NAME
    p_c.font.size = Pt(14)
    p_c.font.bold = True
    p_c.font.color.rgb = COLOR_TEXT_WHITE
    
    # タイトル
    tb_t = slide2.shapes.add_textbox(left + Inches(1.1), top + Inches(0.25), Inches(4.3), Inches(0.5))
    tf_t = tb_t.text_frame
    tf_t.word_wrap = True
    p_t = tf_t.paragraphs[0]
    p_t.text = feat["title"]
    p_t.font.name = FONT_NAME
    p_t.font.size = Pt(16)
    p_t.font.bold = True
    p_t.font.color.rgb = COLOR_TEXT_MAIN
    
    # 説明文
    tb_d = slide2.shapes.add_textbox(left + Inches(0.3), top + Inches(1.1), Inches(5.0), Inches(1.1))
    tf_d = tb_d.text_frame
    tf_d.word_wrap = True
    p_d = tf_d.paragraphs[0]
    p_d.text = feat["desc"]
    p_d.font.name = FONT_NAME
    p_d.font.size = Pt(12)
    p_d.font.color.rgb = COLOR_TEXT_MUTED


# ==========================================
# Slide 3: 全体業務フローとモジュール構成
# ==========================================
slide3 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(slide3, COLOR_BG_LIGHT)
add_header(slide3, "全体業務フローとシステムモジュール構造", "WORKFLOW & ARCHITECTURE")

steps = [
    {"step": "STEP 1", "title": "マスタ管理", "sub": "荷主・商品・倉庫・運送会社・運賃表の設定"},
    {"step": "STEP 2", "title": "入庫処理", "sub": "入庫予定登録 ➔ 確定 ➔ ケース在庫自動生成"},
    {"step": "STEP 3", "title": "倉庫・在庫管理", "sub": "ケース/バラ箱管理 ➔ 残数追跡 ➔ サマリー集計"},
    {"step": "STEP 4", "title": "出庫指示取込", "sub": "Excelインポート ➔ 最安倉庫選定 ➔ FIFO引当"},
    {"step": "STEP 5", "title": "運賃比較・調整", "sub": "複数キャリア運賃比較 ➔ 出荷倉庫確定"},
    {"step": "STEP 6", "title": "送り状出力", "sub": "各社CSVデータ出力 ➔ 配送ラベル発行・出庫確定"}
]

for i, st in enumerate(steps):
    left = Inches(0.6 + i * 2.05)
    top = Inches(1.8)
    
    card = add_card(slide3, left, top, Inches(1.9), Inches(4.5))
    
    # ヘッダー
    hdr = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(1.9), Inches(0.6))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = COLOR_PRIMARY if i % 2 == 0 else COLOR_PRIMARY_DARK
    hdr.line.fill.background()
    p_h = hdr.text_frame.paragraphs[0]
    p_h.text = st["step"]
    p_h.alignment = PP_ALIGN.CENTER
    p_h.font.name = FONT_NAME
    p_h.font.size = Pt(12)
    p_h.font.bold = True
    p_h.font.color.rgb = COLOR_TEXT_WHITE
    
    # タイトル
    tb_t = slide3.shapes.add_textbox(left + Inches(0.1), top + Inches(0.8), Inches(1.7), Inches(0.8))
    tf_t = tb_t.text_frame
    tf_t.word_wrap = True
    p_t = tf_t.paragraphs[0]
    p_t.text = st["title"]
    p_t.alignment = PP_ALIGN.CENTER
    p_t.font.name = FONT_NAME
    p_t.font.size = Pt(14)
    p_t.font.bold = True
    p_t.font.color.rgb = COLOR_TEXT_MAIN
    
    # 説明
    tb_s = slide3.shapes.add_textbox(left + Inches(0.1), top + Inches(1.8), Inches(1.7), Inches(2.4))
    tf_s = tb_s.text_frame
    tf_s.word_wrap = True
    p_s = tf_s.paragraphs[0]
    p_s.text = st["sub"]
    p_s.font.name = FONT_NAME
    p_s.font.size = Pt(11)
    p_s.font.color.rgb = COLOR_TEXT_MUTED


# 共通ヘルパー: 左カード＋右キャプチャレイアウト
def create_feature_slide(prs, title, category, points, img_path):
    slide = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(slide, COLOR_BG_LIGHT)
    add_header(slide, title, category)

    # 左側解説カード
    left_card = add_card(slide, Inches(0.8), Inches(1.4), Inches(4.5), Inches(5.6))
    
    # カード内タイトル
    tb_ct = slide.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(4.1), Inches(0.5))
    tf_ct = tb_ct.text_frame
    p_ct = tf_ct.paragraphs[0]
    p_ct.text = "機能のポイント"
    p_ct.font.name = FONT_NAME
    p_ct.font.size = Pt(16)
    p_ct.font.bold = True
    p_ct.font.color.rgb = COLOR_PRIMARY_DARK

    # 箇条書きポイント
    tb_pts = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(4.1), Inches(4.6))
    tf_pts = tb_pts.text_frame
    tf_pts.word_wrap = True
    
    for i, pt in enumerate(points):
        p = tf_pts.add_paragraph() if i > 0 else tf_pts.paragraphs[0]
        p.text = f"■ {pt['title']}"
        p.font.name = FONT_NAME
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_MAIN
        p.space_before = Pt(8) if i > 0 else Pt(0)

        p_sub = tf_pts.add_paragraph()
        p_sub.text = pt['desc']
        p_sub.font.name = FONT_NAME
        p_sub.font.size = Pt(11)
        p_sub.font.color.rgb = COLOR_TEXT_MUTED
        p_sub.space_after = Pt(6)

    # 右側スクリーンショット
    add_image_with_border(slide, img_path, Inches(5.5), Inches(1.4), Inches(7.0), Inches(5.6))
    return slide


# ==========================================
# Slide 4: 倉庫管理① 在庫一覧
# ==========================================
create_feature_slide(
    prs,
    "倉庫管理：在庫一覧（明細照会・バラ箱残数管理）",
    "INVENTORY MANAGEMENT",
    [
        {"title": "未開封ケース ＆ バラ箱の識別", "desc": "保管ケースが未開封箱か開封済のバラ箱かをリアルタイムに自動識別してバッジ表示します。"},
        {"title": "残存部数の正確な管理", "desc": "バラ箱内の残り部数 / 入数を数値で表示し、現場の端数部数を視覚的かつ正確に把握できます。"},
        {"title": "高精度な絞り込み検索", "desc": "荷主・倉庫・商品コード・ステータス（1:在庫, 11:引当, 21:出庫）の各種複合条件で即座に絞り込み可能です。"}
    ],
    "docs/images/inventory_list.png"
)

# ==========================================
# Slide 5: 倉庫管理② 在庫数サマリー
# ==========================================
create_feature_slide(
    prs,
    "倉庫管理：在庫数サマリー（リアルタイム集計照会）",
    "INVENTORY MANAGEMENT",
    [
        {"title": "倉庫×荷主×商品別の自動集計", "desc": "保管中（有効）な在庫データをリアルタイムにグループ集計し、拠点ごとの総ケース数と部数を一覧化します。"},
        {"title": "利用可能在庫と引当済数の可視化", "desc": "即時出荷可能な「利用可能在庫箱数・部数」と、出荷手配でキープされた「引当済ケース数」を区別表示します。"},
        {"title": "マルチ倉庫・マルチ荷主対応", "desc": "複数拠点・複数荷主の在庫状況を一画面でまとめて把握でき、倉庫間の在庫移動判断を支援します。"}
    ],
    "docs/images/inventory_summary.png"
)

# ==========================================
# Slide 6: 入庫管理（入庫予定・入庫登録）
# ==========================================
slide6 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(slide6, COLOR_BG_LIGHT)
add_header(slide6, "入庫管理：入庫予定一覧 ＆ 入庫登録", "INBOUND MANAGEMENT")

# 左側解説カード
add_card(slide6, Inches(0.8), Inches(1.4), Inches(4.5), Inches(5.6))
tb_s6 = slide6.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(4.1), Inches(5.2))
tf_s6 = tb_s6.text_frame
tf_s6.word_wrap = True

pts6 = [
    {"title": "入庫予定の一元管理", "desc": "入荷予定日時、予定ケース数、荷主・倉庫情報を一覧で把握。作業予定の平準化をサポートします。"},
    {"title": "スムーズな入庫確定処理", "desc": "現場での受入検品後、一括または個別で入庫実績を登録し即座に在庫ステータスへ反映します。"},
    {"title": "ケース在庫自動生成", "desc": "入庫確定時に指定ケース数分の在庫レコード（未開封箱 is_loose=false）を自動生成し、迅速に出荷可能化します。"}
]

for i, pt in enumerate(pts6):
    p = tf_s6.add_paragraph() if i > 0 else tf_s6.paragraphs[0]
    p.text = f"■ {pt['title']}"
    p.font.name = FONT_NAME
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    p.space_before = Pt(8) if i > 0 else Pt(0)

    p_sub = tf_s6.add_paragraph()
    p_sub.text = pt['desc']
    p_sub.font.name = FONT_NAME
    p_sub.font.size = Pt(11)
    p_sub.font.color.rgb = COLOR_TEXT_MUTED
    p_sub.space_after = Pt(6)

# 右側2枚のスクリーンショット
add_image_with_border(slide6, "docs/images/inbound_list.png", Inches(5.5), Inches(1.4), Inches(7.0), Inches(2.65))
add_image_with_border(slide6, "docs/images/inbound_register.png", Inches(5.5), Inches(4.25), Inches(7.0), Inches(2.75))


# ==========================================
# Slide 7: 出庫管理① 出荷指示読込＆自動引当
# ==========================================
create_feature_slide(
    prs,
    "出庫管理：出荷指示読込 ＆ 自動倉庫選定・引当",
    "OUTBOUND MANAGEMENT",
    [
        {"title": "Excelファイルのワンクリック取込", "desc": "出荷先や指示部数が記載されたExcel（.xlsx）を取り込むだけで、自動で出荷指示グループを発行・処理します。"},
        {"title": "最安出荷倉庫の自動選定", "desc": "お届け先郵便番号と指示部数・重量に基づき、最適な出荷倉庫および最安配送ルートを全自動計算して決定します。"},
        {"title": "FIFO × バラ優先の高度な引当ロジック", "desc": "最古の既存バラ箱を残さず優先引き当てし、不足分のみ未開封箱を開封する高精度な自動在庫キープを行います。"}
    ],
    "docs/images/outbound_import_instruction.png"
)

# ==========================================
# Slide 8: 出庫管理② 出庫記録一覧＆進捗管理
# ==========================================
create_feature_slide(
    prs,
    "出庫管理：出庫記録一覧 ＆ リアルタイム進捗管理",
    "OUTBOUND MANAGEMENT",
    [
        {"title": "出荷進捗の一元ビジュアル管理", "desc": "確認中・予定・送状出力・請求済・該当料金無し・在庫切れの各種ステータスをバッジで一目で識別できます。"},
        {"title": "ワンクリック画面連携", "desc": "「確認中」や「該当料金無し」バッジを押下することで、対応が必要な「運賃比較」や「倉庫確定」画面へ即座に移動可能です。"},
        {"title": "出荷グループ別フィルタリング", "desc": "バッチ単位で発番されたグループコードにより、該当ファイル全体の出荷進行状況を直感的に追跡できます。"}
    ],
    "docs/images/outbound_list.png"
)

# ==========================================
# Slide 9: 出庫管理③ 出荷倉庫確定
# ==========================================
create_feature_slide(
    prs,
    "出庫管理：出荷倉庫確定 ＆ 手動調整機能",
    "OUTBOUND MANAGEMENT",
    [
        {"title": "自動選定結果の確認と最適調整", "desc": "システムが自動判定した最安出荷倉庫と適用料金を確認し、必要に応じて配送条件を個別に調整できます。"},
        {"title": "手動変更時の在庫再引当連動", "desc": "出荷倉庫を手動で変更した場合でも、バックエンドで即座に該当倉庫の在庫状況を再検索し、FIFO引当を再計算します。"},
        {"title": "例外ステータスのクリア", "desc": "在庫切れや該当料金無しのレコードについて、倉庫変更や条件調整を行うことで出荷可能ステータスへ修正できます。"}
    ],
    "docs/images/outbound_confirm_warehouse.png"
)

# ==========================================
# Slide 10: 出庫管理④ 高度な運賃比較＆コスト最適化
# ==========================================
create_feature_slide(
    prs,
    "出庫管理：高度な運賃比較 ＆ コスト最適化エンジン",
    "OUTBOUND MANAGEMENT",
    [
        {"title": "複数キャリア運賃の即時シミュレーション", "desc": "ヤマト運輸・佐川急便・日本郵便など、対応運送会社の料金表をリアルタイムで横断比較・最安値を強調表示します。"},
        {"title": "多様な運賃テーブルの網羅", "desc": "サイズ・重量・都道府県・距離別運賃・倉庫個別の割増率・個別運賃設定まで考慮した厳密な金額計算を行います。"},
        {"title": "物流コストの大幅な削減", "desc": "出荷データごとに最も安価な配送事業者・配送方法（路線/個配）を選択でき、年間配送費の直接削減に貢献します。"}
    ],
    "docs/images/outbound_compare_rates.png"
)

# ==========================================
# Slide 11: 出庫管理⑤ 送り状出力＆配送ラベル発行
# ==========================================
create_feature_slide(
    prs,
    "出庫管理：送り状データCSV出力 ＆ 配送ラベル発行",
    "OUTBOUND MANAGEMENT",
    [
        {"title": "主要配送ソフト用CSVの自動生成", "desc": "確定した出庫データから、ヤマトB2 Cloud、佐川e飛伝等に直接取り込める形式の送り状CSVをワンクリックで生成します。"},
        {"title": "配送ラベルの即時発行・印刷", "desc": "倉庫現場で必要な配送ラベルや納品指示書を画面から直接プレビュー＆プリンター発行できます。"},
        {"title": "出庫確定とステータス連動", "desc": "送り状データ出力・ラベル発行と同時に在庫を正式に「21:出庫済」に更新し、二重出庫を防止します。"}
    ],
    "docs/images/outbound_shipping_label.png"
)

# ==========================================
# Slide 12: マスタ管理 物流＆運賃計算基盤
# ==========================================
slide12 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(slide12, COLOR_BG_LIGHT)
add_header(slide12, "マスタ管理：物流オペレーション ＆ 運賃計算の基盤設定", "MASTER MANAGEMENT")

# 左側解説カード
add_card(slide12, Inches(0.8), Inches(1.4), Inches(4.5), Inches(5.6))
tb_s12 = slide12.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(4.1), Inches(5.2))
tf_s12 = tb_s12.text_frame
tf_s12.word_wrap = True

pts12 = [
    {"title": "基本物流マスタ群", "desc": "荷主、商品、倉庫、運送会社、出荷区分、集荷エリア、郵便番号マスタを統合管理。柔軟なマルチ荷主構築を支援。"},
    {"title": "高機能運賃設定エンジン", "desc": "運賃表マスタ、距離マスタ、距離別運賃、倉庫距離割増率、個別運賃マスタにより、複雑な計算ルールを正確に再現。"},
    {"title": "運用変更への高い拡張性", "desc": "運送会社の運賃改定や新倉庫開設時も、マスタ設定の更新のみで即座に全自動選定エンジンに反映可能です。"}
]

for i, pt in enumerate(pts12):
    p = tf_s12.add_paragraph() if i > 0 else tf_s12.paragraphs[0]
    p.text = f"■ {pt['title']}"
    p.font.name = FONT_NAME
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_MAIN
    p.space_before = Pt(8) if i > 0 else Pt(0)

    p_sub = tf_s12.add_paragraph()
    p_sub.text = pt['desc']
    p_sub.font.name = FONT_NAME
    p_sub.font.size = Pt(11)
    p_sub.font.color.rgb = COLOR_TEXT_MUTED
    p_sub.space_after = Pt(6)

# 右側2枚のスクリーンショット
add_image_with_border(slide12, "docs/images/master_freight_table.png", Inches(5.5), Inches(1.4), Inches(7.0), Inches(2.65))
add_image_with_border(slide12, "docs/images/master_product.png", Inches(5.5), Inches(4.25), Inches(7.0), Inches(2.75))


# ==========================================
# Slide 13: 今後の実装課題① WMS基盤機能の拡張
# ==========================================
slide13 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(slide13, COLOR_BG_LIGHT)
add_header(slide13, "今後の実装課題①：WMS基盤機能の拡張（未実装機能の洗い出し）", "FUTURE ISSUES & ROADMAP")

unimplemented_features = [
    {
        "title": "請求・保管料計算機能 (Billing)",
        "icon": "￥",
        "color": COLOR_PRIMARY,
        "items": [
            "荷主別の坪貸し / パレット / ケース保管料の自動集計",
            "入出庫作業費・梱包資材費・ラベル貼り作業費の計算",
            "荷主締め日処理 ＆ 月次請求書データ（PDF/CSV）発行"
        ]
    },
    {
        "title": "作業報告・レポート出力 (Reporting)",
        "icon": "📊",
        "color": COLOR_ACCENT_BLUE,
        "items": [
            "日次/月次の荷主別・倉庫別入出庫作業実績レポート",
            "在庫回転率・長期間滞留在庫の警告アラート出力",
            "運賃シミュレーションによる運賃削減成果ダッシュボード"
        ]
    },
    {
        "title": "棚卸・棚差調整機能 (Stock Taking)",
        "icon": "📋",
        "color": COLOR_ACCENT_AMBER,
        "items": [
            "定期/臨時棚卸のカウント実査データ入力画面",
            "理論在庫 vs 実在庫の差異（棚差）の自動検出と承認",
            "棚卸実施中の該当倉庫ステータス制御 ＆ 監査トラッキング"
        ]
    },
    {
        "title": "詳細ロケーション管理 (Location)",
        "icon": "📍",
        "color": COLOR_ACCENT_EMERALD,
        "items": [
            "エリア・ラック・棚・段単位のロケーションアドレス指定",
            "フリーロケーション運用 ＆ 入出庫作業動線の最適化",
            "商品属性（重量・出荷頻度）に応じた最適ロケ配置提案"
        ]
    }
]

for i, feat in enumerate(unimplemented_features):
    col = i % 2
    row = i // 2
    left = Inches(0.8 + col * 5.9)
    top = Inches(1.5 + row * 2.7)
    
    add_card(slide13, left, top, Inches(5.6), Inches(2.4))
    
    # ヘッダータイトル
    hdr_box = slide13.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(5.6), Inches(0.5))
    hdr_box.fill.solid()
    hdr_box.fill.fore_color.rgb = feat["color"]
    hdr_box.line.fill.background()
    p_h = hdr_box.text_frame.paragraphs[0]
    p_h.text = f"{feat['icon']}  {feat['title']}"
    p_h.font.name = FONT_NAME
    p_h.font.size = Pt(14)
    p_h.font.bold = True
    p_h.font.color.rgb = COLOR_TEXT_WHITE
    
    # リスト項目
    tb_l = slide13.shapes.add_textbox(left + Inches(0.2), top + Inches(0.6), Inches(5.2), Inches(1.7))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    for j, item in enumerate(feat["items"]):
        p = tf_l.add_paragraph() if j > 0 else tf_l.paragraphs[0]
        p.text = f"• {item}"
        p.font.name = FONT_NAME
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_TEXT_MAIN
        p.space_after = Pt(4)


# ==========================================
# Slide 14: 今後の実装課題② 荷主・作業者連携＆将来拡張性
# ==========================================
slide14 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(slide14, COLOR_BG_LIGHT)
add_header(slide14, "今後の実装課題②：荷主・倉庫作業者連携 ＆ 将来拡張性", "FUTURE ISSUES & EXTENSION")

ext_topics = [
    {
        "title": "荷主向けWebポータル\n（荷主マイページ連携）",
        "tag": "荷主連携",
        "color": COLOR_PRIMARY,
        "points": [
            "荷主が自社のリアルタイム在庫数や出荷進捗を直接閲覧",
            "配送コスト削減実績シミュレーションの荷主への可視化",
            "荷主からの出荷指示データWeb直接登録・承認フロー",
            "【課題】マルチテナント型権限管理とデータの厳格な隔離"
        ]
    },
    {
        "title": "現場作業者向けデジタル化\n（ハンディ・HT連携）",
        "tag": "作業者連携",
        "color": COLOR_ACCENT_BLUE,
        "points": [
            "バーコード / QRスキャンによる検品・ペーパーレスピッキング",
            "ハンディターミナル(HT)やスマホアプリからの即時データ更新",
            "誤検品・誤出荷のリアルタイムアラート検知機能",
            "【課題】現場の通信環境に依存しないオフライン処理対応"
        ]
    },
    {
        "title": "外部システム自動連動\n（API / EDI連携基盤）",
        "tag": "外部連携",
        "color": COLOR_ACCENT_EMERALD,
        "points": [
            "Shopify / ネクストエンジン等のECカート/モール自動連動",
            "基幹システム(ERP)や会計ソフトへの売上・請求データ自動送信",
            "各運送会社荷物追跡APIとの自動ステータス同期",
            "【課題】外部API仕様変更に対応する柔軟なアダプタ層設計"
        ]
    }
]

for i, topic in enumerate(ext_topics):
    left = Inches(0.8 + i * 3.95)
    top = Inches(1.5)
    
    card = add_card(slide14, left, top, Inches(3.8), Inches(5.4))
    
    # タグ
    tag_b = slide14.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.3), top + Inches(0.3), Inches(1.2), Inches(0.35))
    tag_b.fill.solid()
    tag_b.fill.fore_color.rgb = topic["color"]
    tag_b.line.fill.background()
    p_t = tag_b.text_frame.paragraphs[0]
    p_t.text = topic["tag"]
    p_t.alignment = PP_ALIGN.CENTER
    p_t.font.name = FONT_NAME
    p_t.font.size = Pt(10)
    p_t.font.bold = True
    p_t.font.color.rgb = COLOR_TEXT_WHITE
    
    # タイトル
    tb_t = slide14.shapes.add_textbox(left + Inches(0.2), top + Inches(0.75), Inches(3.4), Inches(0.9))
    tf_t = tb_t.text_frame
    tf_t.word_wrap = True
    p_title = tf_t.paragraphs[0]
    p_title.text = topic["title"]
    p_title.font.name = FONT_NAME
    p_title.font.size = Pt(15)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_TEXT_MAIN
    
    # 区切り線
    line = slide14.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.3), top + Inches(1.7), Inches(3.2), Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_CARD_BORDER
    line.line.fill.background()
    
    # ポイント
    tb_p = slide14.shapes.add_textbox(left + Inches(0.2), top + Inches(1.8), Inches(3.4), Inches(3.4))
    tf_p = tb_p.text_frame
    tf_p.word_wrap = True
    for j, pt in enumerate(topic["points"]):
        p = tf_p.add_paragraph() if j > 0 else tf_p.paragraphs[0]
        p.text = f"✔ {pt}"
        p.font.name = FONT_NAME
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_PRIMARY_DARK if "【課題】" in pt else COLOR_TEXT_MAIN
        if "【課題】" in pt:
            p.font.bold = True
        p.space_after = Pt(6)


# ==========================================
# Slide 15: 今後のロードマップと開発推進課題
# ==========================================
slide15 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(slide15, COLOR_BG_LIGHT)
add_header(slide15, "今後の開発ロードマップ ＆ システム推進上の課題", "ROADMAP & NEXT STEPS")

phases = [
    {
        "phase": "PHASE 1 (現状)",
        "sub": "基本機能 ＆ 運賃最適化",
        "color": COLOR_PRIMARY,
        "items": ["入出庫基本処理", "バラ箱/残数在庫管理", "自動倉庫選定", "高度運賃比較・ラベル発行"]
    },
    {
        "phase": "PHASE 2 (短期課題)",
        "sub": "WMS基盤機能拡張",
        "color": COLOR_ACCENT_BLUE,
        "items": ["請求・保管料自動計算", "作業報告/帳票出力", "棚卸・棚差調整機能", "ロケーション管理強化"]
    },
    {
        "phase": "PHASE 3 (中長期)",
        "sub": "エコシステム連携・デジタル化",
        "color": COLOR_ACCENT_EMERALD,
        "items": ["荷主向けWebポータル", "ハンディ(HT)検品連携", "EC/ERP自動API連携", "AI需給・在庫予測"]
    }
]

for i, ph in enumerate(phases):
    left = Inches(0.8 + i * 3.95)
    top = Inches(1.5)
    
    card = add_card(slide15, left, top, Inches(3.8), Inches(3.2))
    
    # ヘッダー
    hdr = slide15.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(3.8), Inches(0.7))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = ph["color"]
    hdr.line.fill.background()
    p_h = hdr.text_frame.paragraphs[0]
    p_h.text = ph["phase"]
    p_h.alignment = PP_ALIGN.CENTER
    p_h.font.name = FONT_NAME
    p_h.font.size = Pt(14)
    p_h.font.bold = True
    p_h.font.color.rgb = COLOR_TEXT_WHITE
    
    p_s = hdr.text_frame.add_paragraph()
    p_s.text = ph["sub"]
    p_s.alignment = PP_ALIGN.CENTER
    p_s.font.name = FONT_NAME
    p_s.font.size = Pt(10)
    p_s.font.color.rgb = COLOR_TEXT_WHITE
    
    # 項目リスト
    tb_l = slide15.shapes.add_textbox(left + Inches(0.2), top + Inches(0.8), Inches(3.4), Inches(2.3))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    for j, item in enumerate(ph["items"]):
        p = tf_l.add_paragraph() if j > 0 else tf_l.paragraphs[0]
        p.text = f"▶ {item}"
        p.font.name = FONT_NAME
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_TEXT_MAIN
        p.space_after = Pt(4)

# 下部：システム推進課題
add_card(slide15, Inches(0.8), Inches(4.9), Inches(11.7), Inches(2.1), bg_color=RGBColor(241, 245, 249))

tb_ch = slide15.shapes.add_textbox(Inches(1.0), Inches(5.0), Inches(11.3), Inches(1.9))
tf_ch = tb_ch.text_frame
tf_ch.word_wrap = True

p_cht = tf_ch.paragraphs[0]
p_cht.text = "🔑 実装・運用推進に向けた3大技術検討課題"
p_cht.font.name = FONT_NAME
p_cht.font.size = Pt(14)
p_cht.font.bold = True
p_cht.font.color.rgb = COLOR_PRIMARY_DARK

key_issues = [
    "【1. 計算パフォーマンス最適化】: 荷主・出荷件数増大に伴う運賃マスタ・在庫マスタのクロス集計計算の高速化・非同期バッチ処理化。",
    "【2. データセキュリティ ＆ マルチテナント】: 荷主ポータル公開時における荷主間データの完全論理隔離とアクセス制限。",
    "【3. 現場UI/UX ＆ オペレーション適合】: 倉庫現場作業者が迷わず操作できるタブレット/ハンディUIと誤操作防止ガード。"
]

for issue in key_issues:
    p = tf_ch.add_paragraph()
    p.text = issue
    p.font.name = FONT_NAME
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT_MAIN
    p.space_before = Pt(3)


# ==========================================
# Slide 16: まとめ
# ==========================================
slide16 = prs.slides.add_slide(blank_slide_layout)
set_slide_bg(slide16, COLOR_BG_DARK)

tb_m1 = slide16.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(11), Inches(1.0))
tf_m1 = tb_m1.text_frame
p_m1 = tf_m1.paragraphs[0]
p_m1.text = "まとめ：YK特化WMSがもたらす価値と今後の進化"
p_m1.font.name = FONT_NAME
p_m1.font.size = Pt(28)
p_m1.font.bold = True
p_m1.font.color.rgb = COLOR_TEXT_WHITE

cards_summary = [
    {
        "title": "現状提供価値（強み）",
        "desc": "✔ バラ箱残数管理による現場在庫の精度向上\n✔ 自動倉庫選定 ＆ 高度運賃比較による配送コスト最適化\n✔ Excel指示取込から送り状発行までの一気通貫オペレーション"
    },
    {
        "title": "今後の進化（ロードマップ）",
        "desc": "✔ 請求・保管料自動計算 ＆ 業務報告レポートによる管理業務効率化\n✔ 荷主向けWebポータル ＆ 現場ハンディ(HT)連携によるDX推進\n✔ 外部API基盤拡張によるシームレスなエコシステム構築"
    }
]

for i, cs in enumerate(cards_summary):
    left = Inches(1.2 + i * 5.6)
    top = Inches(2.8)
    
    card = slide16.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.3), Inches(3.2))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(30, 41, 59)
    card.line.color.rgb = COLOR_PRIMARY if i == 0 else COLOR_ACCENT_BLUE
    card.line.width = Pt(2)
    
    tb_ct = slide16.shapes.add_textbox(left + Inches(0.3), top + Inches(0.3), Inches(4.7), Inches(0.5))
    p_ct = tb_ct.text_frame.paragraphs[0]
    p_ct.text = cs["title"]
    p_ct.font.name = FONT_NAME
    p_ct.font.size = Pt(18)
    p_ct.font.bold = True
    p_ct.font.color.rgb = COLOR_ACCENT_BLUE if i == 0 else COLOR_ACCENT_EMERALD
    
    tb_cd = slide16.shapes.add_textbox(left + Inches(0.3), top + Inches(0.9), Inches(4.7), Inches(2.0))
    tf_cd = tb_cd.text_frame
    tf_cd.word_wrap = True
    p_cd = tf_cd.paragraphs[0]
    p_cd.text = cs["desc"]
    p_cd.font.name = FONT_NAME
    p_cd.font.size = Pt(13)
    p_cd.font.color.rgb = COLOR_TEXT_WHITE

# フッター
tb_end = slide16.shapes.add_textbox(Inches(1.2), Inches(6.3), Inches(11), Inches(0.5))
p_end = tb_end.text_frame.paragraphs[0]
p_end.text = "ご清聴ありがとうございました。"
p_end.alignment = PP_ALIGN.CENTER
p_end.font.name = FONT_NAME
p_end.font.size = Pt(16)
p_end.font.bold = True
p_end.font.color.rgb = COLOR_TEXT_MUTED


# PowerPoint保存
output_path = "YK特化WMS_システム概要プレゼンテーション.pptx"
prs.save(output_path)
print(f"Successfully generated PowerPoint file: {os.path.abspath(output_path)}")
